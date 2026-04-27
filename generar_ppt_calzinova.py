#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Genera PPTX: propuesta técnica + manual de usuario (sin tabla de presupuestos)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ASSETS = (
    "/Users/luigiarmandoroblespalacios/.cursor/projects/"
    "Users-luigiarmandoroblespalacios-Desktop-taekwondo-erp/assets"
)

OUT = (
    "/Users/luigiarmandoroblespalacios/Desktop/taekwondo-erp/"
    "Calzinova_Propuesta_Tecnica_y_Manual_Usuario.pptx"
)

# Tema: Strategic Growth
C_PRIMARY = RGBColor(0x00, 0x4D, 0x40)
C_GOLD = RGBColor(0xC5, 0xB3, 0x58)
C_SLATE = RGBColor(0x70, 0x80, 0x90)
C_SMOKE = RGBColor(0xF5, 0xF5, 0xF5)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def style_title(shape, size=32):
    if not shape.has_text_frame:
        return
    p = shape.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.font.name = "Helvetica Neue"


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left = Inches(0.6)
    top = Inches(1.0)
    box = slide.shapes.add_textbox(left, top, Inches(8.5), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.font.name = "Helvetica Neue"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = C_SLATE
        p2.space_after = Pt(6)
    return slide


def add_bullet_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title and content
    style_title(slide.shapes.title, 28)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(15)
        p.font.name = "Helvetica Neue"
        p.font.color.rgb = C_SLATE
        p.space_after = Pt(4)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_table_slide(prs, title, headers, rows, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    style_title(slide.shapes.title, 24)
    slide.shapes.title.text = title
    rows_n = 1 + len(rows)
    cols = len(headers)
    left = Inches(0.4)
    top = Inches(1.2)
    width = Inches(9.0)
    height = Inches(0.2 * rows_n + 0.3)
    table = slide.shapes.add_table(rows_n, cols, left, top, width, min(height, Inches(5.0))).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = C_PRIMARY
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            p = table.cell(r, c).text_frame.paragraphs[0]
            p.text = str(val)[:200]
            p.font.size = Pt(9)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_image_slide(prs, title, relpath, foot="", img_height=Inches(4.2)):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide.shapes.title, 22)
    slide.shapes.title.text = title
    import os
    path = os.path.join(ASSETS, relpath)
    if not os.path.isfile(path):
        raise FileNotFoundError(path)
    # Centrar imagen
    w_slide = 10.0
    left = (Inches(10) - Inches(8.0)) / 2 + Inches(0.3)  # ~1.0
    left = Inches(0.5)
    top = Inches(1.1)
    pic = slide.shapes.add_picture(path, left, top, width=Inches(8.0), height=img_height)
    if foot:
        footb = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(8.5), Inches(0.5))
        ft = footb.text_frame
        ft.text = foot
        ft.paragraphs[0].font.size = Pt(11)
        ft.paragraphs[0].font.italic = True
        ft.paragraphs[0].font.color.rgb = C_SLATE
    slide.notes_slide.notes_text_frame.text = (
        f"Imagen: {relpath}. {foot}"
    )
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Portada
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s0.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.8), Inches(2.2))
    t = tb.text_frame
    t.text = "Proyecto de Ingeniería de Software"
    t.paragraphs[0].font.size = Pt(18)
    t.paragraphs[0].font.color.rgb = C_SLATE
    t.add_paragraph().text = "CALZINOVA"
    t.paragraphs[1].font.size = Pt(40)
    t.paragraphs[1].font.bold = True
    t.paragraphs[1].font.color.rgb = C_PRIMARY
    t.add_paragraph().text = (
        "Marketplace y panel de gestión de indicadores comerciales"
    )
    t.paragraphs[2].font.size = Pt(16)
    t.paragraphs[2].font.color.rgb = C_SLATE
    t.add_paragraph().text = (
        "Propuesta técnica (Calzado Libertad) + manual de usuario"
    )
    t.paragraphs[3].font.size = Pt(14)
    t.paragraphs[3].font.color.rgb = C_GOLD
    t.add_paragraph().text = "Responsables: Luigi Armando Robles Palacios, Jhayron Joseph Quiliche Valladolid"
    t.paragraphs[4].font.size = Pt(12)
    t.paragraphs[4].font.color.rgb = C_SLATE
    t.add_paragraph().text = "Trujillo, La Libertad — 2026"
    t.paragraphs[5].font.size = Pt(12)
    t.paragraphs[5].font.color.rgb = C_SLATE

    # 2. Logo
    add_image_slide(
        prs,
        "Identidad de marca: CALZINOVA",
        "PHOTO-2026-04-19-23-57-27-20a3958d-4284-4991-9b56-f5a3d58ab296.png",
        "Icono de calzado y barras de crecimiento: artesanía + datos.",
    )

    # 3. Estructura de la guía
    add_bullet_slide(
        prs,
        "Estructura de esta presentación (guía + manual)",
        [
            "Parte A — Qué pide y qué entrega el proyecto (propuesta técnica, sin cifras de presupuesto).",
            "Parte B — Cómo se organiza la información en el sistema: módulos de MYPE (inventario, producción, ventas, capital humano).",
            "Parte C — Seguridad, usuarios y roles (multi-tenancy).",
            "Parte D — Cómo usar la web pública, el acceso empresarial y el panel (paso a paso con capturas).",
        ],
    )

    # 4. Resumen propuesta
    add_bullet_slide(
        prs,
        "Idea del proyecto (visión operativa y comercial)",
        [
            "Un ERP ligero y un marketplace: la fábrica (MYPE) ve inventario, producción, ventas e indicadores en un panel.",
            "El visitante del marketplace explora el catálogo y contacta por WhatsApp; el empresario concentra márgenes, capacidad y clima en un solo lugar.",
            "Cobertura de referencia: ecosistema productivo (Trujillo, El Porvenir) y enfoque en transformación digital del calzado.",
        ],
    )

    # 5. Requerimientos
    add_bullet_slide(
        prs,
        "1. Requerimientos que la empresa debe proveer (arranque del proyecto)",
        [
            "Identidad gráfica: logotipo (PNG/SVG) y catálogo de modelos actuales.",
            "Inventario: fotos, tallas, colores y precios por modelo.",
            "Estructura de costos: costo de producción por par (para calibrar KPIs).",
            "Canal de venta: número de WhatsApp Business (integración y contacto directo).",
        ],
    )

    # 6. Entregables
    add_bullet_slide(
        prs,
        "2. Entregables del sistema",
        [
            "Dominio .com: registro y configuración por 12 meses (según acuerdo comercial).",
            "Panel administrativo: inventario, costos, métricas e indicadores.",
            "Marketplace web: catálogo interactivo e integración con WhatsApp.",
            "Soporte técnico: garantía de mantenimiento y ajustes (ventana de 3 meses en propuesta).",
        ],
    )

    # 7. Infraestructura (sin tablas de monto; sin tabla de “presupuesto de desarrollo”)
    add_bullet_slide(
        prs,
        "3. Inversión en infraestructura (concepto, sin desglose presupuestal)",
        [
            "Hosting en la nube (plan anual) para alojar panel y marketplace con buen desempeño.",
            "Certificado SSL y dominio .com: tráfico cifrado y presencia en Internet bajo un nombre propio (según el paquete contratado).",
            "Esta sección reemplaza la presentación de tablas con costos, por solicitud: los montos se revisan con el proveedor o equipo al momento del despliegue.",
        ],
    )

    # 8. Seguridad propuesta
    add_bullet_slide(
        prs,
        "5. Seguridad y mantenimiento (nivel de servicio acordado)",
        [
            "Backups periódicos: respaldo de la base de datos (por ejemplo, semanal, según política de hosting).",
            "SSL: cifrado para datos en tránsito (ventas e información de clientes).",
        ],
    )

    # 9. Fórmulas KPI
    add_bullet_slide(
        prs,
        "6. Fórmulas de KPIs de negocio (cómo se interpretan en el panel)",
        [
            "Utilidad neta aprox. en %: ((Venta − Costo) / Venta) × 100",
            "Punto de equilibrio: Gastos / (Precio de venta − costo variable unitario).",
        ],
        notes="Apliques estas fórmulas a los costos reales y precios reales; el software muestra el resultado y tendencias, no reemplaza el criterio del empresario.",
    )

    # 10. Anexos técnicos (referencia a URLs de ejemplo en PDF)
    add_bullet_slide(
        prs,
        "Anexos de la propuesta: vistas de referencia (Calzado Libertad / despliegue propuesto)",
        [
            "Anexo A — Panel: ejemplo de URL de administración: admin.calzadolibertad.pe/dashboard (KPI: ventas del mes, utilidad, stock bajo, gráfica ventas vs costos).",
            "Anexo B — Marketplace: ejemplo de URL pública: calzadolibertad.pe, catálogo con enlace a WhatsApp por producto, diseño responsive para móvil.",
        ],
    )

    # 11. Paleta
    add_image_slide(
        prs,
        "Lineamiento visual: paleta “Strategic Growth”",
        "PHOTO-2026-04-19-23-44-50-70b8bb6c-5178-4fef-af16-5249eb57b899.png",
        "Aplica a fondos, botones, alertas y textos: verde #004D40, oro #C5B358, pizarra #708090, humo #F5F5F5.",
    )

    # --- Módulo inventario
    add_table_slide(
        prs,
        "Base de datos MYPE: Inventario (logística) — qué registro en cada campo",
        ["Campo", "Uso", "Indicador UPRIT (referencia)"],
        [
            ("Material", "Insumo (cuero, suela, hilo, pegamento…)", "Control de insumos"),
            ("Unidad de medida", "Pies, docenas, litros, uds…", "Gestión de compras"),
            ("Stock actual", "Lo que hay hoy en almacén", "R4.6 abastecimiento"),
            ("Stock mínimo", "Umbral de alerta de compra", "Alerta de quiebre"),
            ("Costo unitario", "Precio de compra", "R1.7 valorización"),
            ("Valor total", "Stock × costo; capital inmovilizado", "R1.7 capital de trabajo"),
        ],
        notes="Rellene primero unidades y costos; luego el sistema calcula el valor y alerta cuando el stock baja del mínimo.",
    )

    add_table_slide(
        prs,
        "Base de datos MYPE: Producción (operaciones)",
        ["Campo", "Uso", "Indicador UPRIT"],
        [
            ("Código de lote", "ID del pedido/batch (trazabilidad)", "Trazabilidad"),
            ("Modelo", "Diseño o referencia de calzado", "Catálogo"),
            ("Cant. programada", "Pares pedidos", "Meta de producción"),
            ("Cant. ejecutada", "Pares terminados", "R4.2 eficiencia"),
            ("Estado actual", "Corte → Aparado → Armado → Terminado", "R4.4 lead time"),
            ("Fecha entrega", "compromiso con el cliente", "R4.1 cumplimiento"),
        ],
    )

    add_table_slide(
        prs,
        "Base de datos MYPE: Ventas (comercial)",
        ["Campo", "Uso", "Indicador UPRIT"],
        [
            ("Cliente", "Persona o empresa", "Cartera"),
            ("Total venta", "Monto cobrado", "Ingresos brutos"),
            ("Método de pago", "Efectivo, Yape, Plin, transferencia", "R1.1 bancarización"),
            ("Estado de pago", "Pagado / pendiente", "Control de deuda"),
            ("Ticket promedio", "Gasto promedio del cliente en el mes", "R3.3"),
        ],
    )

    add_table_slide(
        prs,
        "Base de datos MYPE: Capital humano (personas y bienestar)",
        ["Campo", "Uso", "Indicador UPRIT"],
        [
            ("Operario", "Nombre", "Registro de personal"),
            ("Especialidad", "Cortador, aparador…", "Especialización"),
            ("Horas capacitación", "Aprendizaje continuo", "H2.1 capacitación"),
            ("Incidentes", "Riesgos o accidentes", "H5.7 seguridad"),
        ],
    )

    # Seguridad PDF 2
    add_bullet_slide(
        prs,
        "Base de datos: seguridad y usuarios (objetivo)",
        [
            "Multi-tenancy: cada taller (MYPE) ve solo sus propios datos; se evita mezclar información entre fábricas.",
            "Identidad, correo, rol y contraseña aseguran quién hace qué, y dejan rastro (auditoría básica).",
        ],
    )

    add_table_slide(
        prs,
        "Módulo de identidad: campos y propósito",
        ["Campo (sistema)", "Nombre para el usuario", "Propósito en el negocio"],
        [
            ("Name", "Nombre visible", "Auditoría: quién operó el sistema"),
            ("Email", "Credencial de acceso", "Login e identificador único; notificaciones"),
            ("Role", "Nivel de permisos", "Qué menús y acciones ve cada persona"),
            ("Password", "Contraseña (cifrada)", "Nadie, ni el admin, ve la clave en claro: solo restablecimiento/rotación controlada"),
        ],
    )

    add_table_slide(
        prs,
        "Definición de roles (explicar en planta: quién es quién en CALZINOVA)",
        ["Rol", "Alcance del acceso"],
        [
            (
                "Admin (soporte)",
                "Plataforma completa, nuevas MYPEs, salud de servidor y tareas de soporte.",
            ),
            (
                "MYPE (dueño de taller)",
                "Control total de su inventario, producción, ventas y clima. No ve datos de otros talleres.",
            ),
            (
                "Cliente (comprador)",
                "Uso restringido: favoritos, carrito o listas; sin acceso a datos operativos de terceras empresas.",
            ),
        ],
    )

    # --- Manual: pantallas
    add_image_slide(
        prs,
        "Paso 1: Acceso empresarial (inicio de sesión)",
        "afa70704-962b-4919-954d-a36c5eb1645b-730f6254-8727-4217-8d49-0f8968123546.png",
        "Correo, contraseña, recordar sesión, recuperación, registro, indicador de conexión segura.",
        img_height=Inches(3.8),
    )

    add_bullet_slide(
        prs,
        "Cómo usar la pantalla de inicio (para el final del entrenamiento)",
        [
            "Escriba su correo empresarial y su contraseña. Use el ícono del ojo para comprobar que escribió bien, sin dejar el teclado sin supervisionar.",
            "“Recordarme” solo en equipos de confianza. “¿Olvidaste tu contraseña?” abre el flujo de restablecimiento (según implementación de correo).",
            "“Regístrate aquí” conduce al alta (si aplica a su fábrica). El aviso de conexión segura (candado) confirma cifrado SSL.",
        ],
    )

    add_image_slide(
        prs,
        "Sitio público: landing (valor y llamadas a la acción) — variante 1",
        "d19b7cf6-aa5d-48b8-b1c3-68e82681ebac-97a797b0-4441-4083-bbe1-019c0bdc1503.png",
        foot="Héroe, “dos pilares” (rentabilidad y capital humano), vistazo a KPI, CTA a marketplace y acceso.",
        img_height=Inches(3.1),
    )

    add_image_slide(
        prs,
        "Sitio público: indicadores científicos y KPIs — variante 2",
        "116be844-9f3a-40c6-bb0f-4ddb8d2e7983-bcd58599-97a1-49f2-abe9-e98bf43ff7e9.png",
        foot="Refuerza los pilares: 56 indicadores financieros y 42 de capital humano; “cerebro del sistema”.",
        img_height=Inches(3.1),
    )

    add_bullet_slide(
        prs,
        "Dónde hacer clic (visitante o empresario)",
        [
            "“Explorar marketplace”: catálogo, filtros y contacto (WhatsApp) por producto.",
            "“Acceso empresarial”: vuelve al login del panel. “Comenzar ahora” impulsa el onboarding.",
        ],
    )

    add_image_slide(
        prs,
        "Marketplace: catálogo, filtros y contacto (WhatsApp)",
        "68fd6300-706c-4aab-9b34-0a00fc8697fc-f24a1dd1-b013-416f-a65b-30813145b115.png",
        "Categoría, talla, rango de precio; en cada tarjeta, precio en soles e ícono WhatsApp.",
    )

    add_bullet_slide(
        prs,
        "Manual: uso del marketplace (comprador)",
        [
            "Use el buscador o la sección (Novedades, Marcas, Ofertas) según esté en su diseño publicado.",
            "Ajuste filtros laterales: categoría, talla EU, precio mínimo y máximo.",
            "Cada ficha resume precio, fabricante o MYPE, y atajo a WhatsApp para concretar pedido, talla y envío por chat.",
        ],
    )

    # Panel: tres capturas
    add_image_slide(
        prs,
        "Panel — Control de producción: resumen (KPIs) y + Nueva Orden",
        "53ea3c38-d299-464a-96a4-342518f852de-8330e27a-0482-46cd-b7d1-afca6f80223f.png",
        "Menú: Dashboard, Inventario, Control de producción, Ventas, Capital humano, Cerrar sesión. Estado vacío: cree la primera orden.",
    )

    add_image_slide(
        prs,
        "Panel — Formulario: crear una orden (lote, modelo, pares, etapa, fecha)",
        "386bf9b7-13ba-464b-99db-86d151213e99-beb5df02-3239-4429-911a-12b8dd20dc3f.png",
        "Filtre por Corte, Aparado, Armado, Terminado. Luego “Crear orden” o “Cancelar”.",
    )

    add_image_slide(
        prs,
        "Panel — Formulario (detalle) y área de listado (cuando aún no hay órdenes)",
        "718cb772-249d-41b4-b8b9-0c8bfce9a19d-5b0be950-3ef2-42ab-90ba-2533a0168bec.png",
        "Mismo flujo: KPI en cero al inicio; al registrar lotes, las tarjetas y la tabla se llenan.",
    )

    add_bullet_slide(
        prs,
        "Qué poner en “Control de producción” (en orden lógico)",
        [
            "1) Haga clic en + Nueva Orden o abra el formulario de creación.",
            "2) Escriba un código de lote único (p. ej. LOTE-2025-001) y el modelo (p. ej. “Zapatilla Deportiva X1”).",
            "3) Indique pares, estado inicial (suele ser Corte al arrancar) y fecha de inicio.",
            "4) Filtre por etapa (Todas, Corte, Aparado, Armado, Terminado) para ver solo lo que le interesa.",
        ],
    )

    add_bullet_slide(
        prs,
        "Cierre: qué explicar a su equipo (recordatorio de la propuesta)",
        [
            "Requisitos, entregables, infraestructura, seguridad y fórmulas: ya cubiertas en diapositivas anteriores (sin tabla de presupuestos, por solicitud).",
            "La aplicación vincula datos de MYPE, marketplace y acceso por roles: úselo de forma consciente; actualice estados reales (producción, pago) para confiar en los indicadores.",
        ],
    )

    s_end = prs.slides.add_slide(prs.slide_layouts[5])
    s_end.shapes.title.text = "Gracias"
    style_title(s_end.shapes.title, 32)
    note = s_end.notes_slide.notes_text_frame
    note.text = (
        "Documento: propuesta técnica (Calzado Libertad) + anexos de datos MYPE y seguridad, "
        "más recorrido de manual con todas las imágenes de interfaz. "
        "Excluida: tabla de presupuesto de desarrollo (sección 4 del PDF de propuesta)."
    )

    prs.save(OUT)
    print("Guardado:", OUT)


if __name__ == "__main__":
    main()
